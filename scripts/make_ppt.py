"""
The Code Room — Capstone PPT generator (v2)
흰 배경 + 어두운 텍스트 + 교육 서사 흐름
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color palette (white bg, dark text) ──────────────────────────────────────
BG       = RGBColor(0xFF, 0xFF, 0xFF)   # white
DARK     = RGBColor(0x11, 0x18, 0x27)   # near-black (main text)
MID      = RGBColor(0x37, 0x41, 0x51)   # dark gray (body text)
DIM      = RGBColor(0x9C, 0xA3, 0xAF)   # muted gray (captions)
NAVY     = RGBColor(0x1E, 0x3A, 0x5F)   # navy accent (titles, kickers)
BLUE     = RGBColor(0x29, 0x80, 0xB9)   # medium blue (highlights)
CARD_BG  = RGBColor(0xF0, 0xF6, 0xFF)   # very light blue (card fill)
RULE_CLR = RGBColor(0xD1, 0xD5, 0xDB)   # light gray (dividers)
GREEN    = RGBColor(0x0D, 0x6E, 0x47)   # forest green (positive callout)
AMBER    = RGBColor(0x92, 0x40, 0x00)   # amber (warning/nuance)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def box(slide, text, x, y, w, h,
        size=20, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def multiline(slide, lines, x, y, w, h, default_size=19, default_color=MID):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col, sz = item, False, default_color, default_size
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else default_color
            sz   = item[3] if len(item) > 3 else default_size
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = col
    return tb

def bullets(slide, items, x, y, w, h, size=19, icon="▸"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            text, bold, col = item, False, MID
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else MID
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6)
        r1 = p.add_run()
        r1.text = icon + "  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = BLUE
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(size)
        r2.font.bold = bold
        r2.font.color.rgb = col
    return tb

def rect(slide, x, y, w, h, fill=CARD_BG, border=RULE_CLR, border_pt=0.75):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(border_pt)
    return sh

def hbar(slide, x, y, w, color=RULE_CLR, pt=0.75):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0))
    sh.line.color.rgb = color
    sh.line.width = Pt(pt)

def left_accent(slide, x, y, h, color=NAVY, w=0.07):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()

def kicker(slide, text, x=0.5, y=0.32):
    box(slide, text, x=x, y=y, w=5, h=0.3, size=11, bold=True, color=BLUE)

def slide_num(slide, n, total=8):
    box(slide, f"{n}  /  {total}", x=12.3, y=7.1, w=1, h=0.35,
        size=11, color=DIM, align=PP_ALIGN.RIGHT)

def top_rule(slide):
    hbar(slide, 0.5, 0.7, 12.3, color=RULE_CLR, pt=0.75)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()

# Decorative: faint dot grid (light gray)
for xi in [1.0, 3.0, 5.5, 8.0, 10.5, 12.5]:
    for yi in [1.2, 2.8, 4.4, 6.0]:
        dot = s.shapes.add_shape(1, Inches(xi), Inches(yi), Inches(0.05), Inches(0.05))
        dot.fill.solid(); dot.fill.fore_color.rgb = RULE_CLR
        dot.line.fill.background()

# Left navy bar accent
bar = s.shapes.add_shape(1, Inches(1.1), Inches(2.0), Inches(0.1), Inches(2.8))
bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()

# Title
box(s, "THE CODE ROOM",
    x=1.4, y=1.95, w=10.5, h=1.5, size=62, bold=True, color=DARK)

# Subtitle
box(s, "방탈출 구조를 활용한 Python 개념 복습 게임",
    x=1.4, y=3.55, w=10.5, h=0.7, size=22, color=NAVY)

# Separator
hbar(s, 1.4, 4.5, 9.5, color=RULE_CLR)

# Meta
box(s, "2026-1 캡스톤 프로젝트   |   김민혁 · 공원호",
    x=1.4, y=4.7, w=9.5, h=0.5, size=16, color=DIM)

slide_num(s, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 문제의식
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "PROBLEM")
top_rule(s)
box(s, "Python 수업은 많아졌지만,\n기억에 잘 안 남는다",
    x=0.5, y=0.75, w=12.3, h=1.3, size=32, bold=True, color=DARK)

# Context block
rect(s, 0.5, 2.2, 12.3, 1.05, fill=CARD_BG, border=RULE_CLR)
box(s, "디지털새싹 · 정보영재 프로그램 · SW·AI 캠프 · 학교 정보 수업에서 Python은 이제 입문 언어로 자리잡고 있다.",
    x=0.7, y=2.32, w=11.9, h=0.8, size=17, color=MID)

# Arrow flow
flow_items = ["개념 설명", "예제 코드 따라치기", "문제 몇 개 풀이", "다음 개념으로"]
for i, text in enumerate(flow_items):
    bx = 0.5 + i * 3.0
    rect(s, bx, 3.5, 2.7, 0.62, fill=BG, border=NAVY, border_pt=1.0)
    box(s, text, x=bx+0.05, y=3.55, w=2.6, h=0.5, size=15, bold=True,
        color=NAVY, align=PP_ALIGN.CENTER)
    if i < 3:
        box(s, "→", x=bx+2.72, y=3.55, w=0.3, h=0.5, size=18, color=DIM,
            align=PP_ALIGN.CENTER)

box(s, "이 방식은 진도를 나가기에는 효율적이지만, 학생 입장에서 다음 문제가 생긴다.",
    x=0.5, y=4.35, w=12.3, h=0.5, size=16, color=MID)

bullets(s, [
    "수업 중에는 따라오지만, 시간이 지나면 기억이 잘 안 남음",
    "문제를 풀 때 '왜 그렇게 푸는지'보다 '예제랑 비슷하게' 푸는 경우가 많음",
    "단순 반복 문제풀이가 이어지면 집중도가 떨어짐",
], x=0.5, y=4.85, w=12.3, h=1.8, size=17)

box(s, "한국과학창의재단 (디지털새싹).   교육부 (2022). 2022 개정 정보과 교육과정 — 컴퓨팅 사고력·실생활 문제 해결·재미있는 프로그래밍 강조.",
    x=0.5, y=6.95, w=12.3, h=0.4, size=10, color=DIM, italic=True)
slide_num(s, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 제작 동기
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "MOTIVATION")
top_rule(s)
box(s, "그냥 푸는 Python 문제를,\n탈출해야 하는 문제로 바꾸면?",
    x=0.5, y=0.75, w=12.3, h=1.3, size=32, bold=True, color=DARK)

# Quote
rect(s, 0.5, 2.2, 12.3, 0.95, fill=CARD_BG, border=NAVY, border_pt=1.2)
box(s, '"문제 1번을 푸세요"  →  "이 코드를 풀어야 다음 방으로 갈 수 있다"',
    x=0.7, y=2.32, w=11.9, h=0.7, size=19, bold=True, color=NAVY, italic=True)

# Table header
rect(s, 0.5, 3.38, 5.9, 0.42, fill=NAVY, border=NAVY)
rect(s, 6.42, 3.38, 6.38, 0.42, fill=NAVY, border=NAVY)
box(s, "방탈출 구조", x=0.5, y=3.38, w=5.9, h=0.42, size=15, bold=True,
    color=BG, align=PP_ALIGN.CENTER)
box(s, "Python 문제 해결", x=6.42, y=3.38, w=6.38, h=0.42, size=15, bold=True,
    color=BG, align=PP_ALIGN.CENTER)

rows = [
    ("단서를 찾는다",        "문제 조건을 읽는다"),
    ("규칙을 추론한다",      "필요한 문법과 로직을 고른다"),
    ("자물쇠를 연다",        "코드를 실행해 정답을 확인한다"),
    ("실패하면 다시 시도",   "오류를 보고 디버깅한다"),
    ("다음 방으로 이동",     "다음 개념으로 확장한다"),
]
row_colors = [BG, CARD_BG, BG, CARD_BG, BG]
for i, (left, right) in enumerate(rows):
    by = 3.82 + i * 0.48
    rect(s, 0.5,  by, 5.9,  0.46, fill=row_colors[i], border=RULE_CLR, border_pt=0.5)
    rect(s, 6.42, by, 6.38, 0.46, fill=row_colors[i], border=RULE_CLR, border_pt=0.5)
    box(s, left,  x=0.65,  y=by+0.05, w=5.6,  h=0.38, size=15, color=MID)
    box(s, right, x=6.57,  y=by+0.05, w=6.1,  h=0.38, size=15, color=MID)
    # center ≈
    if i == 2:
        box(s, "≈", x=6.13, y=by+0.02, w=0.3, h=0.42, size=18, bold=True,
            color=BLUE, align=PP_ALIGN.CENTER)

slide_num(s, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 교육적 근거
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "EDUCATIONAL BASIS")
top_rule(s)
box(s, "학습자는 설명을 듣는 것보다,\n직접 써볼 때 더 잘 배운다",
    x=0.5, y=0.75, w=12.3, h=1.3, size=32, bold=True, color=DARK)

box(s, "The Code Room은 강의를 대체하려는 것이 아니라, 수업 이후 학생이 배운 개념을 직접 꺼내 쓰는 복습 활동으로 설계됐다.",
    x=0.5, y=2.1, w=12.3, h=0.6, size=16, color=MID)

cards = [
    ("Active Learning (Freeman et al., PNAS 2014)",
     "STEM 분야 메타분석 — 능동학습이 전통 강의 대비\n시험 성과를 유의미하게 향상, 실패율을 낮춤"),
    ("Escape Education (Veldkamp et al., 2020)",
     "교육용 방탈출 체계적 문헌 검토\n핵심: 퍼즐 설계와 학습 목표의 정렬이 효과를 결정"),
    ("디지털 방탈출 (Sánchez, ScienceDirect 2023)",
     "디지털 방탈출이 학습 동기와 내용 복습에 긍정적\n학생들이 학습 과정의 질이 향상됐다고 인식"),
    ("Escape Rooms for Programming (BERA, 2024)",
     "프로그래밍 학습용 방탈출 체계적 문헌 검토\n'코딩 + 방탈출'은 검증된 교육 연구 영역"),
]

for i, (title, body) in enumerate(cards):
    col_i = i % 2
    row_i = i // 2
    bx = 0.5 + col_i * 6.45
    by = 2.85 + row_i * 2.3
    rect(s, bx, by, 6.15, 2.1, fill=CARD_BG, border=RULE_CLR)
    left_accent(s, bx, by, 2.1, color=NAVY)
    box(s, title, x=bx+0.2, y=by+0.12, w=5.85, h=0.52, size=13, bold=True, color=NAVY)
    box(s, body,  x=bx+0.2, y=by+0.68, w=5.85, h=1.3,  size=13, color=MID)

slide_num(s, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 설계 방향
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "GAME DESIGN")
top_rule(s)
box(s, "핵심은 게임성이 아니라\n학습 목표와 퍼즐의 연결이다",
    x=0.5, y=0.75, w=12.3, h=1.3, size=32, bold=True, color=DARK)

# Room table header
rect(s, 0.5, 2.18, 1.8, 0.42, fill=NAVY, border=NAVY)
rect(s, 2.32, 2.18, 3.2, 0.42, fill=NAVY, border=NAVY)
rect(s, 5.54, 2.18, 7.24, 0.42, fill=NAVY, border=NAVY)
box(s, "방", x=0.5, y=2.18, w=1.8, h=0.42, size=14, bold=True,
    color=BG, align=PP_ALIGN.CENTER)
box(s, "학습 목표", x=2.32, y=2.18, w=3.2, h=0.42, size=14, bold=True,
    color=BG, align=PP_ALIGN.CENTER)
box(s, "활동", x=5.54, y=2.18, w=7.24, h=0.42, size=14, bold=True,
    color=BG, align=PP_ALIGN.CENTER)

rooms = [
    ("Room 0\n잠긴 서재",   "변수 · 연산\n인덱싱 · 슬라이싱",    "기본 코드 해석, 값 계산, 문자열 자르기"),
    ("Room 1\n신호실",      "문자열 메서드\n조건문 (if / else)", "단서에서 조건을 찾아 코드 작성"),
    ("Room 2\n기록실",      "반복문 (for · while)\n자료구조 (리스트·딕셔너리)", "여러 데이터를 순회하며 정답 도출"),
]
row_colors = [BG, CARD_BG, BG]
for i, (room, goal, act) in enumerate(rooms):
    by = 2.62 + i * 0.88
    rect(s, 0.5,  by, 1.8,  0.86, fill=row_colors[i], border=RULE_CLR, border_pt=0.5)
    rect(s, 2.32, by, 3.2,  0.86, fill=row_colors[i], border=RULE_CLR, border_pt=0.5)
    rect(s, 5.54, by, 7.24, 0.86, fill=row_colors[i], border=RULE_CLR, border_pt=0.5)
    box(s, room, x=0.55,  y=by+0.06, w=1.7,  h=0.75, size=13, bold=True, color=NAVY)
    box(s, goal, x=2.37,  y=by+0.06, w=3.1,  h=0.75, size=13, color=MID)
    box(s, act,  x=5.59,  y=by+0.06, w=7.14, h=0.75, size=13, color=MID)

# Bonus note
rect(s, 0.5, 5.26, 12.28, 0.52, fill=RGBColor(0xFF,0xF8,0xF0),
     border=RGBColor(0xD4,0xA0,0x17), border_pt=0.75)
box(s, "★ 보너스 룸 (선택): 3개 방 클리어 후 진입 가능한 시크릿 방 — 역방향 슬라이싱 심화 문제 1개, 별도 스토리 엔딩",
    x=0.7, y=5.3, w=12.0, h=0.42, size=13, color=AMBER, bold=False)

# Design principles
box(s, "설계 원칙", x=0.5, y=5.95, w=3, h=0.38, size=14, bold=True, color=NAVY)
hbar(s, 0.5, 6.33, 12.3, color=RULE_CLR)
principles = [
    "문제를 풀어야만 다음 방으로 이동 — 강제적 능동 참여",
    "단서는 방 오브젝트처럼 제시 — 문제 조건을 자연스럽게 읽게 유도",
    "요구 문법 사용 여부를 AST로 검사 — 정답만이 아니라 과정도 검증",
    "실패 시 감점 없음, 재시도 무제한 — 부담 없는 시행착오",
]
# 2x2 grid for principles
for i, text in enumerate(principles):
    col_i = i % 2
    row_i = i // 2
    bx = 0.5 + col_i * 6.3
    by = 6.4 + row_i * 0.45
    left_accent(s, bx, by+0.05, 0.32, color=BLUE, w=0.05)
    box(s, text, x=bx+0.15, y=by, w=6.0, h=0.42, size=13, color=MID)

slide_num(s, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 주요 게임 요소 / 학생 경험 흐름
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "STUDENT EXPERIENCE")
top_rule(s)
box(s, "학생 경험 흐름",
    x=0.5, y=0.75, w=12.3, h=0.9, size=34, bold=True, color=DARK)

# Flow steps
steps = [
    ("방 입장",      "새로운 학습 단위 시작"),
    ("단서 탐색",    "오브젝트 클릭으로\n문제 조건 확인"),
    ("코드 작성",    "Python 에디터에\n직접 코드 입력"),
    ("실행 · 확인",  "즉각 피드백\n(정답 / 오류 메시지)"),
    ("오류 수정",    "디버깅 후 재시도"),
    ("자물쇠 해제",  "6개 단서 코드 조합\n→ 도어 키패드 입력"),
    ("다음 방",      "학습 진행에 대한\n보상 + 새 개념"),
]

for i, (title, desc) in enumerate(steps):
    bx = 0.35 + i * 1.84
    rect(s, bx, 1.95, 1.68, 1.8, fill=CARD_BG, border=NAVY, border_pt=1.0)
    # colored top strip
    top_strip = s.shapes.add_shape(1, Inches(bx), Inches(1.95), Inches(1.68), Inches(0.18))
    top_strip.fill.solid(); top_strip.fill.fore_color.rgb = NAVY
    top_strip.line.fill.background()
    box(s, title, x=bx+0.06, y=2.17, w=1.56, h=0.42,
        size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    box(s, desc, x=bx+0.06, y=2.6, w=1.56, h=1.1,
        size=12, color=MID, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        box(s, "→", x=bx+1.72, y=2.65, w=0.15, h=0.4, size=14,
            color=DIM, align=PP_ALIGN.CENTER)

# Element table
box(s, "주요 요소", x=0.5, y=4.1, w=3, h=0.38, size=14, bold=True, color=NAVY)
hbar(s, 0.5, 4.48, 12.3, color=RULE_CLR)

elements = [
    ("방",       "하나의 Python 개념 학습 단위 (Room 0–2, 총 3개)"),
    ("단서",     "문제 조건을 방 오브젝트처럼 자연스럽게 제시하는 장치"),
    ("코드 창",  "CodeMirror 에디터 — 문법 강조·자동완성 지원"),
    ("실행 결과","Pyodide로 브라우저 내 Python 직접 실행, 즉각 피드백"),
    ("자물쇠",   "6개 퍼즐을 모두 풀어야 열리는 도어 키패드 — 문제 해결의 목표"),
    ("다음 방",  "방 클리어 = 학습 단위 완료 → 새로운 개념 영역으로 이동"),
]
for i, (el, desc) in enumerate(elements):
    col_i = i % 2
    row_i = i // 2
    bx = 0.5 + col_i * 6.45
    by = 4.55 + row_i * 0.62
    box(s, el,   x=bx+0.02, y=by, w=1.5,  h=0.55, size=14, bold=True, color=NAVY)
    box(s, desc, x=bx+1.55, y=by, w=4.85, h=0.55, size=14, color=MID)

slide_num(s, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 기대 효과
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "EXPECTED OUTCOMES")
top_rule(s)
box(s, "기대 효과",
    x=0.5, y=0.75, w=12.3, h=0.9, size=34, bold=True, color=DARK)

effects = [
    ("집중도 향상",
     "단순 문제풀이보다 '다음 방으로 이동'이라는 목표가 있어 몰입을 유도할 수 있음"),
    ("기억 유지 보조",
     "문법을 개념명으로만 기억하는 것이 아니라, 특정 탈출 상황과 연결해서 기억하게 됨"),
    ("능동적 문제 해결",
     "학생이 조건을 읽고, 필요한 문법을 선택하고, 코드를 직접 수정하는 경험"),
    ("실패-재시도 학습",
     "틀려도 감점이 없고 재시도가 자연스러운 구조 — 시행착오를 통한 개념 내면화"),
    ("교사 활용성",
     "디지털새싹·정보영재·방과후·학교 정보 수업의 Python 복습 활동으로 사용 가능\n(브라우저만 있으면 설치 불필요)"),
]
accent_colors = [NAVY, BLUE, RGBColor(0x0D,0x6E,0x47), NAVY, BLUE]
for i, (title, body) in enumerate(effects):
    by = 1.75 + i * 0.99
    bar = s.shapes.add_shape(1, Inches(0.5), Inches(by+0.08), Inches(0.07), Inches(0.72))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent_colors[i]
    bar.line.fill.background()
    box(s, f"{'①②③④⑤'[i]}  {title}", x=0.72, y=by, w=3.0, h=0.42,
        size=16, bold=True, color=DARK)
    box(s, body, x=0.72, y=by+0.44, w=12.0, h=0.52, size=15, color=MID)

# Nuance note
rect(s, 0.5, 6.7, 12.3, 0.58, fill=RGBColor(0xFF,0xF8,0xF0),
     border=RGBColor(0xD4,0xA0,0x17), border_pt=0.75)
box(s, "⚠  The Code Room은 Python 실력을 단번에 올려주는 도구가 아니라, 이미 배운 개념을 더 능동적으로 다시 꺼내 쓰게 만드는 복습 환경입니다.",
    x=0.68, y=6.76, w=12.0, h=0.45, size=12, color=AMBER)

slide_num(s, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 마무리
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "CONCLUSION")
top_rule(s)
box(s, "코딩 개념을 써서 탈출하는 경험",
    x=0.5, y=0.75, w=12.3, h=0.9, size=34, bold=True, color=DARK)

# Main message box
rect(s, 0.5, 1.88, 12.3, 1.5,
     fill=RGBColor(0xEF,0xF6,0xFF), border=NAVY, border_pt=1.5)
left_accent(s, 0.5, 1.88, 1.5, color=NAVY, w=0.12)
box(s, '"Python 문제를 더 많이 푸는 것이 아니라,\nPython 문제를 더 몰입해서 풀 수 있는 상황을 만든다."',
    x=0.78, y=1.96, w=11.9, h=1.3, size=20, bold=True, color=NAVY, italic=True)

bullets(s, [
    ("The Code Room = 코딩을 외우는 활동이 아니라, 코딩 개념을 사용해 탈출하는 경험", True, DARK),
    ("방탈출 구조는 프로그래밍 문제 해결의 사고 흐름(단서→추론→실행→재시도)과 자연스럽게 맞물림", False, MID),
    ("교육 연구 기반 설계: Active learning + 교육용 방탈출 + Serious game 원칙 반영", False, MID),
], x=0.5, y=3.55, w=12.3, h=2.0, size=18)

# Final line
rect(s, 0.5, 5.6, 12.3, 0.82, fill=CARD_BG, border=RULE_CLR)
box(s, '"학생은 문제를 풀어야 다음 방으로 갈 수 있고, 그 과정에서 Python 개념은 자연스럽게 다시 사용된다."',
    x=0.65, y=5.68, w=12.0, h=0.65, size=16, color=NAVY, italic=True, align=PP_ALIGN.CENTER)

# References
hbar(s, 0.5, 6.58, 12.3, color=RULE_CLR)
box(s, "References", x=0.5, y=6.63, w=2.5, h=0.3, size=11, bold=True, color=NAVY)
refs_text = (
    "교육부 (2022). 2022 개정 교육과정 정보과.   "
    "Freeman et al. (2014). Active learning increases student performance. PNAS.   "
    "Veldkamp et al. (2020). Escape education: systematic review. ScienceDirect.   "
    "Sánchez (2023). Digital educational escape room. ScienceDirect.   "
    "Wouters et al. (2019). Serious games meta-analysis. TU/e.   "
    "Burak et al. (2024). Escape rooms for learning programming. BERA."
)
box(s, refs_text, x=0.5, y=6.95, w=12.3, h=0.42, size=9, color=DIM, italic=True)

slide_num(s, 8)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = r"C:\Users\user\Desktop\ETDR\The_Code_Room_Capstone.pptx"
prs.save(out_path)
print(f"Saved → {out_path}")
